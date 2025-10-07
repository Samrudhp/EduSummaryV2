"""
Text extraction utilities for PDF, PPT, and DOCX files
"""
import re
from typing import List, Dict
import PyPDF2
import pdfplumber
from pptx import Presentation
from docx import Document


def clean_text(text: str) -> str:
    """Clean and normalize extracted text"""
    # Remove excessive whitespace
    text = re.sub(r'\s+', ' ', text)
    # Remove special characters but keep punctuation
    text = re.sub(r'[^\w\s.,;:!?\-\(\)\[\]\'\"]+', '', text)
    # Normalize line breaks
    text = text.replace('\n', ' ').replace('\r', ' ')
    return text.strip()


def extract_from_pdf(file_path: str) -> str:
    """Extract text from PDF using pdfplumber with page preservation"""
    text = ""
    try:
        with pdfplumber.open(file_path) as pdf:
            for page_num, page in enumerate(pdf.pages, 1):
                page_text = page.extract_text()
                if page_text:
                    # Keep page markers for section detection
                    text += f"[PAGE_{page_num}]\n{page_text}\n\n"
    except Exception as e:
        print(f"Error extracting PDF with pdfplumber: {e}")
        # Fallback to PyPDF2
        try:
            with open(file_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                for page_num, page in enumerate(pdf_reader.pages, 1):
                    page_text = page.extract_text()
                    if page_text:
                        text += f"[PAGE_{page_num}]\n{page_text}\n\n"
        except Exception as e2:
            print(f"Error extracting PDF with PyPDF2: {e2}")
            raise
    
    return text  # Don't clean yet - we need structure for section detection


def extract_from_pptx(file_path: str) -> str:
    """Extract text from PowerPoint"""
    text = ""
    try:
        prs = Presentation(file_path)
        for slide_num, slide in enumerate(prs.slides, 1):
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += f"\n{shape.text}"
            text += "\n\n"
    except Exception as e:
        print(f"Error extracting PPTX: {e}")
        raise
    
    return clean_text(text)


def extract_from_docx(file_path: str) -> str:
    """Extract text from Word document"""
    text = ""
    try:
        doc = Document(file_path)
        for para in doc.paragraphs:
            if para.text.strip():
                text += para.text + "\n"
        text += "\n"
    except Exception as e:
        print(f"Error extracting DOCX: {e}")
        raise
    
    return clean_text(text)


def extract_text(file_path: str, file_type: str) -> str:
    """Main extraction function"""
    if file_type == "pdf":
        return extract_from_pdf(file_path)
    elif file_type == "pptx":
        return extract_from_pptx(file_path)
    elif file_type == "docx":
        return extract_from_docx(file_path)
    else:
        raise ValueError(f"Unsupported file type: {file_type}")


def extract_sections(text: str) -> List[Dict[str, str]]:
    """
    1. Document structure analysis (headings, formatting)
    2. Semantic coherence (topic shifts)
    3. Statistical features (line length, capitalization patterns)
    4. Academic paper structure detection
    """
    print("\n" + "="*60)
    print("🔍 ANALYZING DOCUMENT WITH TECHNIQUES...")
    print("="*60)
    
    sections = []
    lines = text.split('\n')
    
    # ========== PHASE 1: TEXT CLEANING & NORMALIZATION ==========
    print("\n[Phase 1] Cleaning and normalizing text...")
    cleaned_lines = []
    page_boundaries = []  # Track page breaks for context
    
    for idx, line in enumerate(lines):
        line = line.strip()
        
        # Track page boundaries
        if line.startswith('[PAGE_'):
            page_match = re.match(r'\[PAGE_(\d+)\]', line)
            if page_match:
                page_boundaries.append((idx, int(page_match.group(1))))
            continue
        
        # Skip empty or very short lines
        if len(line) < 3:
            continue
            
        # Skip metadata patterns
        if any([
            '@' in line and '.' in line,  # Email addresses
            line.startswith('http'),  # URLs
            line.startswith('www.'),
            line.startswith('doi:'),  # DOI references
            re.match(r'^\d{4}-\d{4}', line),  # ISSN numbers
            re.match(r'^[\d\s\-\(\)]+$', line),  # Phone/reference numbers
            len(re.findall(r'[^a-zA-Z0-9\s]', line)) > len(line) * 0.4  # Too many special chars
        ]):
            continue
        
        cleaned_lines.append({
            'text': line,
            'original_idx': idx,
            'length': len(line),
            'words': len(line.split()),
            'is_upper': line.isupper(),
            'is_title_case': line.istitle(),
            'starts_with_capital': line[0].isupper() if line else False,
            'has_numbers': bool(re.search(r'\d', line))
        })
    
    print(f"  ✓ Cleaned: {len(lines)} → {len(cleaned_lines)} lines")
    
    # ========== PHASE 2: HEADING DETECTION WITH SCORING ==========
    print("\n[Phase 2] Detecting headings with confidence scoring...")
    heading_candidates = []
    
    # Academic section patterns (highest confidence)
    ACADEMIC_SECTIONS = {
        r'^abstract$': ('Abstract', 10),
        r'^introduction$': ('Introduction', 10),
        r'^(related work|literature review|background)$': ('Related Work', 10),
        r'^(methodology|methods?)$': ('Methodology', 10),
        r'^(experiments?|experimental setup)$': ('Experiments', 10),
        r'^(results?|findings?)$': ('Results', 10),
        r'^(discussion|analysis)$': ('Discussion', 10),
        r'^(conclusion|conclusions?)$': ('Conclusion', 10),
        r'^(references?|bibliography)$': ('References', 10),
        r'^(acknowledgements?|acknowledgments?)$': ('Acknowledgements', 10),
        r'^(appendix|appendices)$': ('Appendix', 10),
    }
    
    for idx, line_obj in enumerate(cleaned_lines):
        line = line_obj['text']
        score = 0
        heading_type = 'unknown'
        cleaned_title = line
        
        # Check academic sections first
        for pattern, (title, conf) in ACADEMIC_SECTIONS.items():
            if re.match(pattern, line.lower()):
                score = conf
                heading_type = 'academic_section'
                cleaned_title = title
                break
        
        if score == 0:  # Not an academic section, check other patterns
            
            # Pattern 1: Numbered sections (Chapter 1, Section 2.1, etc.)
            if re.match(r'^(chapter|section|part|unit|module|lesson|article)\s+\d+', line, re.IGNORECASE):
                score = 9
                heading_type = 'numbered_section'
                cleaned_title = re.sub(r'^(chapter|section|part|unit|module|lesson|article)\s+', '', line, flags=re.IGNORECASE).strip()
                cleaned_title = f"{line.split()[0].title()} {cleaned_title}"
            
            # Pattern 2: Dotted numbering (1. Introduction, 2.3 Methods)
            elif re.match(r'^\d+\.(\d+\.?)?\s+[A-Z][a-zA-Z\s]{2,}$', line):
                score = 8
                heading_type = 'dotted_number'
                # Remove number prefix
                cleaned_title = re.sub(r'^\d+\.(\d+\.?)?\s+', '', line)
            
            # Pattern 3: ALL CAPS multi-word headings
            elif (line_obj['is_upper'] and 
                  line_obj['words'] >= 2 and 
                  8 <= line_obj['length'] <= 60 and
                  not line_obj['has_numbers']):
                score = 7
                heading_type = 'all_caps'
                cleaned_title = line.title()  # Convert to title case
            
            # Pattern 4: Title Case without ending punctuation
            elif (line_obj['is_title_case'] and
                  line_obj['words'] >= 2 and
                  10 <= line_obj['length'] <= 80 and
                  not line.endswith(('.', ',', ';', ':', '?', '!'))):
                score = 6
                heading_type = 'title_case'
            
            # Pattern 5: Short capitalized lines (potential headings)
            elif (line_obj['starts_with_capital'] and
                  2 <= line_obj['words'] <= 8 and
                  15 <= line_obj['length'] <= 70 and
                  not line.endswith(('.', ','))):
                score = 5
                heading_type = 'short_capitalized'
        
        if score >= 5:  # Only keep candidates with decent confidence
            heading_candidates.append({
                'idx': idx,
                'text': cleaned_title,
                'original_text': line,
                'type': heading_type,
                'score': score,
                'line_obj': line_obj
            })
    
    print(f"  ✓ Found {len(heading_candidates)} heading candidates")
    
    # ========== PHASE 3: HEADING VALIDATION & FILTERING ==========
    print("\n[Phase 3] Validating and filtering headings...")
    
    # Sort by position
    heading_candidates.sort(key=lambda x: x['idx'])
    
    # Remove duplicates and too-close headings
    validated_headings = []
    for i, heading in enumerate(heading_candidates):
        # Skip if too close to previous heading (< 3 lines)
        if validated_headings and heading['idx'] - validated_headings[-1]['idx'] < 3:
            # Keep the one with higher score
            if heading['score'] > validated_headings[-1]['score']:
                validated_headings[-1] = heading
            continue
        
        # Check if next few lines look like content (not another heading)
        has_content = False
        for j in range(heading['idx'] + 1, min(heading['idx'] + 5, len(cleaned_lines))):
            next_line = cleaned_lines[j]
            if next_line['length'] > 50 and not next_line['is_upper']:
                has_content = True
                break
        
        if has_content or heading['score'] >= 9:  # Academic sections don't need content check
            validated_headings.append(heading)
    
    print(f"  ✓ Validated: {len(validated_headings)} high-quality headings")
    
    # ========== PHASE 4: SECTION CONSTRUCTION ==========
    print("\n[Phase 4] Building sections with content...")
    
    if len(validated_headings) >= 2:
        for i, heading in enumerate(validated_headings):
            start_idx = heading['idx']
            end_idx = validated_headings[i + 1]['idx'] if i + 1 < len(validated_headings) else len(cleaned_lines)
            
            # Extract content between headings
            section_lines = cleaned_lines[start_idx + 1:end_idx]
            content = ' '.join([l['text'] for l in section_lines])
            
            # Only include sections with substantial content
            if len(content) > 100:
                preview = content[:250] + "..." if len(content) > 250 else content
                
                sections.append({
                    "id": f"section_{i}",
                    "title": heading['text'][:100],
                    "preview": preview,
                    "content": content,
                    "type": heading['type'],
                    "confidence": heading['score']
                })
                
                print(f"  ✓ Section {i+1}: '{heading['text'][:50]}' ({len(content)} chars, score: {heading['score']})")
    
    # ========== PHASE 5: INTELLIGENT FALLBACK ==========
    if len(sections) < 2:
        print("\n[Phase 5] No clear structure - using intelligent content division...")
        
        # Combine all text
        full_text = ' '.join([l['text'] for l in cleaned_lines])
        total_chars = len(full_text)
        
        # Calculate optimal section count (3-8 sections based on length)
        optimal_sections = min(8, max(3, total_chars // 2000))
        
        # Group lines into paragraphs based on semantic boundaries
        paragraphs = []
        current_para = []
        
        for i, line_obj in enumerate(cleaned_lines):
            line = line_obj['text']
            
            # Start new paragraph on:
            # - Short lines after long content
            # - Capitalized starts after lowercase ends
            # - Topic shifts (simple heuristic)
            if current_para:
                prev_line = current_para[-1]
                if (len(current_para) > 3 and len(line) < 40) or \
                   (len(' '.join(current_para)) > 200 and line[0].isupper()):
                    paragraphs.append(' '.join(current_para))
                    current_para = []
            
            current_para.append(line)
        
        if current_para:
            paragraphs.append(' '.join(current_para))
        
        # Group paragraphs into sections
        if paragraphs:
            paras_per_section = max(1, len(paragraphs) // optimal_sections)
            
            for i in range(optimal_sections):
                start = i * paras_per_section
                end = start + paras_per_section if i < optimal_sections - 1 else len(paragraphs)
                
                if start >= len(paragraphs):
                    break
                
                section_content = ' '.join(paragraphs[start:end])
                
                if len(section_content) > 150:
                    # Extract meaningful title from content
                    sentences = section_content.split('.')
                    title = sentences[0][:80].strip() if sentences else f"Part {i + 1}"
                    
                    # Clean up title
                    title = re.sub(r'^\d+\s+', '', title)  # Remove leading numbers
                    if not title or len(title) < 10:
                        title = f"Section {i + 1}"
                    
                    preview = section_content[:250] + "..." if len(section_content) > 250 else section_content
                    
                    sections.append({
                        "id": f"section_{i}",
                        "title": title,
                        "preview": preview,
                        "content": section_content,
                        "type": "content_based",
                        "confidence": 5
                    })
                    
                    print(f"  ✓ Auto-section {i+1}: '{title[:50]}' ({len(section_content)} chars)")
    
    print(f"\n{'='*60}")
    print(f"✨ EXTRACTED {len(sections)} SECTIONS SUCCESSFULLY!")
    print(f"{'='*60}\n")
    
    return sections


def chunk_text(text: str, chunk_size: int = 300, overlap: int = 30, 
               section_id: str = None, section_title: str = None) -> List[Dict[str, any]]:
    """
    Chunk text into smaller pieces with metadata
    chunk_size: approximate number of tokens per chunk
    overlap: number of tokens to overlap between chunks
    section_id: ID of the section this text belongs to
    section_title: Title of the section
    """
    # Approximate: 1 token ≈ 4 characters
    char_chunk_size = chunk_size * 4
    char_overlap = overlap * 4
    
    chunks = []
    words = text.split()
    
    current_chunk = []
    current_length = 0
    chunk_id = 0
    
    for word in words:
        current_chunk.append(word)
        current_length += len(word) + 1  # +1 for space
        
        if current_length >= char_chunk_size:
            chunk_text = ' '.join(current_chunk)
            chunks.append({
                'chunk_id': chunk_id,
                'text': chunk_text,
                'metadata': {
                    'chunk_id': chunk_id,
                    'char_count': len(chunk_text),
                    'section_id': section_id,
                    'section_title': section_title
                }
            })
            chunk_id += 1
            
            # Keep overlap for next chunk
            overlap_words = int(len(current_chunk) * (char_overlap / char_chunk_size))
            current_chunk = current_chunk[-overlap_words:] if overlap_words > 0 else []
            current_length = sum(len(w) + 1 for w in current_chunk)
    
    # Add remaining chunk
    if current_chunk:
        chunk_text = ' '.join(current_chunk)
        chunks.append({
            'chunk_id': chunk_id,
            'text': chunk_text,
            'metadata': {
                'chunk_id': chunk_id,
                'char_count': len(chunk_text),
                'section_id': section_id,
                'section_title': section_title
            }
        })
    
    return chunks
